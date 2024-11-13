import google.generativeai as genai
import os
from pptx import Presentation
import re
from flask import Flask, render_template, request, send_file
from PIL import Image
from io import BytesIO

# Configure API key
os.environ["API_KEY"] = 'AIzaSyCd_BLtcUii28mpEDGz-vAkKuSE4HgghTI'  # Replace with your actual API key
genai.configure(api_key=os.environ["API_KEY"])

# Initialize the Gemini model at the global level
model = genai.GenerativeModel('gemini-1.5-flash-latest')

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'  # Folder to save uploaded images
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Route for the main page
@app.route('/')
def index():
    return render_template('index.html')

# Function to convert image to a supported format
def convert_image(input_image_path, output_image_path):
    with Image.open(input_image_path) as img:
        img.convert("RGB").save(output_image_path, "JPEG")

# Route to generate the PowerPoint
@app.route('/generate', methods=['POST'])
def generate():
    try:
        topic = request.form['topic']
        custom_background = request.files.get('custom_background')

        # Generate default background image
        background_image = generate_default_background()

        # Handle custom background image upload
        if custom_background:
            custom_image_path = os.path.join(app.config['UPLOAD_FOLDER'], custom_background.filename)
            custom_background.save(custom_image_path)
            background_image = custom_image_path

        # Generate structured content for slides
        response = model.generate_content(f"Create a structured presentation outline for : {topic}. Each slide should contain bullet points summarizing key aspects of the topic.")
        content = response.text.strip()

        # Create a PowerPoint presentation
        prs = Presentation()

        # Split the content into slides
        slides_content = content.split('**Slide')

        for slide_text in slides_content:
            slide_text = slide_text.strip()
            if slide_text:
                current_slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using a title and content layout
                lines = slide_text.split('\n')
                title = lines[0].replace(':', '').strip()
                title = re.sub(r'^\*+|\*+$', '', title)
                title = re.sub(r'\*\*', '', title)
                title = re.sub(r'^\d+\s+', '', title)

                current_slide.shapes.title.text = title
                content_box = current_slide.placeholders[1]
                content_box.text = ''

                for line in lines[1:]:
                    line = line.strip()
                    if line:
                        line = re.sub(r'^\*+|\*+$', '', line)
                        line = re.sub(r'\*\*', '', line)
                        content_box.text += line + '\n'

                # Set background image
                left = top = 0
                if background_image:
                    if background_image.lower().endswith('.webp'):
                        converted_image_path = os.path.join(app.config['UPLOAD_FOLDER'], "converted_image.jpg")  # Temporary file for converted image
                        convert_image(background_image, converted_image_path)
                        pic = current_slide.shapes.add_picture(converted_image_path, left, top, width=prs.slide_width, height=prs.slide_height)
                    else:
                        pic = current_slide.shapes.add_picture(background_image, left, top, width=prs.slide_width, height=prs.slide_height)

                    # Move the picture to the back by removing and re-adding
                    current_slide.shapes._spTree.remove(pic._element)
                    current_slide.shapes._spTree.insert(2, pic._element)  # Adjust the index if needed

        # Save the presentation to a BytesIO object
        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)  # Move to the beginning of the BytesIO object

        return send_file(ppt_io, as_attachment=True, download_name=f"{topic.replace(' ', '_')}_presentation.pptx", mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return f"An error occurred: {str(e)}"

def generate_default_background():
    # Implement your logic to generate a default background image
    default_image_path = os.path.join(app.config['UPLOAD_FOLDER'], "default_background.jpg")
    img = Image.new('RGB', (1280, 720), color='lightblue')  # Example: light blue background
    img.save(default_image_path)
    return default_image_path

if __name__ == '__main__':
    app.run(debug=True)